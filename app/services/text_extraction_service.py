import logging
import boto3
import io
import PyPDF2
import docx
import openpyxl
import pptx
from typing import Optional, Dict, Any
from botocore.exceptions import ClientError
import mimetypes
import chardet
from config.settings import settings

logger = logging.getLogger(__name__)


class TextExtractionService:
    """Service for extracting text content from various file types"""
    
    def __init__(self):
        self.s3_client = boto3.client(
            's3',
            aws_access_key_id=settings.aws_access_key_id,
            aws_secret_access_key=settings.aws_secret_access_key,
            aws_session_token=settings.aws_session_token if hasattr(settings, 'aws_session_token') else None,
            region_name=settings.aws_region
        )
        self.bucket_name = settings.s3_bucket_name
    
    async def extract_text_from_s3(self, s3_key: str, mime_type: str) -> Dict[str, Any]:
        """Extract text content from a file stored in S3"""
        try:
            # Download file from S3
            response = self.s3_client.get_object(Bucket=self.bucket_name, Key=s3_key)
            file_content = response['Body'].read()
            
            # Extract text based on file type
            text = await self._extract_text(file_content, mime_type, s3_key)
            
            # Calculate statistics
            stats = self._calculate_stats(text)
            
            return {
                'text': text,
                'stats': stats,
                'success': True
            }
            
        except Exception as e:
            logger.error(f"Error extracting text from S3 file {s3_key}: {e}")
            return {
                'text': '',
                'stats': {},
                'success': False,
                'error': str(e)
            }
    
    async def _extract_text(self, file_content: bytes, mime_type: str, filename: str) -> str:
        """Extract text based on MIME type"""
        
        # PDF files
        if mime_type == 'application/pdf':
            return self._extract_pdf_text(file_content)
        
        # Word documents
        elif mime_type in ['application/vnd.openxmlformats-officedocument.wordprocessingml.document', 
                          'application/msword']:
            return self._extract_docx_text(file_content)
        
        # Excel files
        elif mime_type in ['application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
                          'application/vnd.ms-excel']:
            return self._extract_excel_text(file_content)
        
        # PowerPoint files
        elif mime_type in ['application/vnd.openxmlformats-officedocument.presentationml.presentation',
                          'application/vnd.ms-powerpoint']:
            return self._extract_pptx_text(file_content)
        
        # Text files
        elif mime_type.startswith('text/') or mime_type in ['application/json', 'application/xml']:
            return self._extract_plain_text(file_content)
        
        # CSV files
        elif mime_type == 'text/csv' or filename.endswith('.csv'):
            return self._extract_csv_text(file_content)
        
        else:
            logger.warning(f"Unsupported file type: {mime_type}")
            return ""
    
    def _extract_pdf_text(self, file_content: bytes) -> str:
        """Extract text from PDF files"""
        try:
            pdf_file = io.BytesIO(file_content)
            pdf_reader = PyPDF2.PdfReader(pdf_file)
            
            text_parts = []
            for page_num in range(len(pdf_reader.pages)):
                page = pdf_reader.pages[page_num]
                text = page.extract_text()
                if text:
                    text_parts.append(text)
            
            return '\n'.join(text_parts)
        except Exception as e:
            logger.error(f"Error extracting PDF text: {e}")
            return ""
    
    def _extract_docx_text(self, file_content: bytes) -> str:
        """Extract text from Word documents"""
        try:
            doc_file = io.BytesIO(file_content)
            doc = docx.Document(doc_file)
            
            text_parts = []
            for paragraph in doc.paragraphs:
                if paragraph.text:
                    text_parts.append(paragraph.text)
            
            # Also extract text from tables
            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        if cell.text:
                            text_parts.append(cell.text)
            
            return '\n'.join(text_parts)
        except Exception as e:
            logger.error(f"Error extracting DOCX text: {e}")
            return ""
    
    def _extract_excel_text(self, file_content: bytes) -> str:
        """Extract text from Excel files"""
        try:
            excel_file = io.BytesIO(file_content)
            workbook = openpyxl.load_workbook(excel_file, read_only=True, data_only=True)
            
            text_parts = []
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                text_parts.append(f"Sheet: {sheet_name}")
                
                for row in sheet.iter_rows(values_only=True):
                    row_text = ' | '.join(str(cell) if cell is not None else '' for cell in row)
                    if row_text.strip():
                        text_parts.append(row_text)
            
            return '\n'.join(text_parts)
        except Exception as e:
            logger.error(f"Error extracting Excel text: {e}")
            return ""
    
    def _extract_pptx_text(self, file_content: bytes) -> str:
        """Extract text from PowerPoint presentations"""
        try:
            pptx_file = io.BytesIO(file_content)
            presentation = pptx.Presentation(pptx_file)
            
            text_parts = []
            for slide_num, slide in enumerate(presentation.slides, 1):
                text_parts.append(f"Slide {slide_num}:")
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        text_parts.append(shape.text)
            
            return '\n'.join(text_parts)
        except Exception as e:
            logger.error(f"Error extracting PPTX text: {e}")
            return ""
    
    def _extract_plain_text(self, file_content: bytes) -> str:
        """Extract text from plain text files"""
        try:
            # Detect encoding
            detected = chardet.detect(file_content)
            encoding = detected.get('encoding', 'utf-8')
            
            # Decode with detected encoding
            text = file_content.decode(encoding, errors='ignore')
            return text
        except Exception as e:
            logger.error(f"Error extracting plain text: {e}")
            # Try with default UTF-8
            try:
                return file_content.decode('utf-8', errors='ignore')
            except:
                return ""
    
    def _extract_csv_text(self, file_content: bytes) -> str:
        """Extract text from CSV files"""
        try:
            import csv
            import io
            
            # Detect encoding
            detected = chardet.detect(file_content)
            encoding = detected.get('encoding', 'utf-8')
            
            # Decode and parse CSV
            text_content = file_content.decode(encoding, errors='ignore')
            csv_reader = csv.reader(io.StringIO(text_content))
            
            text_parts = []
            for row in csv_reader:
                if row:
                    text_parts.append(' | '.join(row))
            
            return '\n'.join(text_parts)
        except Exception as e:
            logger.error(f"Error extracting CSV text: {e}")
            return ""
    
    def _calculate_stats(self, text: str) -> Dict[str, Any]:
        """Calculate text statistics"""
        if not text:
            return {
                'char_count': 0,
                'word_count': 0,
                'line_count': 0,
                'estimated_tokens': 0
            }
        
        # Basic statistics
        char_count = len(text)
        word_count = len(text.split())
        line_count = len(text.splitlines())
        
        # Estimate tokens (rough approximation: 1 token ≈ 4 characters)
        estimated_tokens = char_count // 4
        
        return {
            'char_count': char_count,
            'word_count': word_count,
            'line_count': line_count,
            'estimated_tokens': estimated_tokens
        }


# Create singleton instance
text_extraction_service = TextExtractionService()